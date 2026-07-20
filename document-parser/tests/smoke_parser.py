from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation


def run_parser(binary: Path, source: Path) -> dict:
    result = subprocess.run(
        [str(binary), "--path", str(source)],
        check=False,
        capture_output=True,
        text=True,
        timeout=60,
    )
    if result.returncode != 0:
        raise RuntimeError(result.stderr or result.stdout)
    return json.loads(result.stdout)


def main() -> int:
    binary = Path(sys.argv[1]).resolve()
    with tempfile.TemporaryDirectory(prefix="agnes-parser-smoke-") as directory:
        source = Path(directory) / "smoke.docx"
        document = Document()
        document.add_heading("Smoke test", level=1)
        document.add_paragraph("Frozen parser is operational.")
        document.save(source)
        presentation_path = Path(directory) / "smoke.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Frozen presentation"
        slide.placeholders[1].text = "PowerPoint parsing is operational."
        presentation.save(presentation_path)
        workbook_path = Path(directory) / "smoke.xlsx"
        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Frozen workbook"
        worksheet.append(["status", "value"])
        worksheet.append(["operational", 1])
        workbook.save(workbook_path)

        try:
            payloads = [
                run_parser(binary, source),
                run_parser(binary, presentation_path),
                run_parser(binary, workbook_path),
            ]
        except RuntimeError as error:
            print(error, file=sys.stderr)
            return 1
    assert all(payload["schema_version"] == 1 for payload in payloads)
    assert all(payload["chunks"] for payload in payloads)
    assert "Frozen parser" in payloads[0]["chunks"][0]["content"]
    assert "Frozen presentation" in payloads[1]["chunks"][0]["content"]
    assert "operational" in payloads[2]["chunks"][0]["content"]
    print("document-parserd smoke test passed")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
