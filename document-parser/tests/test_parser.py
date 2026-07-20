import json
import sys
import zipfile
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

import pytest

import app.parser as parser_module
from app.parser import DocumentParseError, PARSER_PROFILE, parse_document
from document_parserd import main as parser_main


def test_docx_preserves_heading_and_table_metadata(tmp_path: Path) -> None:
    path = tmp_path / "report.docx"
    document = Document()
    document.add_heading("实验结果", level=1)
    document.add_paragraph("苹果样本表现稳定。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "名称"
    table.cell(0, 1).text = "数量"
    table.cell(1, 0).text = "苹果"
    table.cell(1, 1).text = "3"
    document.save(path)

    result = parse_document(str(path))

    assert result["parser_profile"] == PARSER_PROFILE
    assert result["media_type"].endswith("wordprocessingml.document")
    assert any(chunk["section_path"] == "实验结果" for chunk in result["chunks"])
    table_chunk = next(chunk for chunk in result["chunks"] if chunk["metadata"]["kind"] == "table")
    assert "苹果" in table_chunk["content"]
    assert table_chunk["page"] is None


def test_office_media_type_hint_supports_extensionless_staged_files(tmp_path: Path) -> None:
    path = tmp_path / "staged-source"
    document = Document()
    document.add_paragraph("网盘暂存文件仍可解析。")
    document.save(path)

    result = parse_document(
        str(path),
        "remote.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

    assert result["title"] == "remote.docx"
    assert "网盘暂存文件" in result["chunks"][0]["content"]


def test_pptx_uses_slide_number_and_title(tmp_path: Path) -> None:
    path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "系统架构"
    slide.placeholders[1].text = "桌面端\n同步服务"
    presentation.save(path)

    result = parse_document(str(path), "演示文稿")

    chunk = next(chunk for chunk in result["chunks"] if chunk["metadata"]["kind"] == "slide")
    assert result["title"] == "演示文稿"
    assert chunk["page"] == 1
    assert chunk["section_path"] == "系统架构"
    assert "桌面端" in chunk["content"]


def test_xlsx_preserves_sheet_and_splits_large_tables(tmp_path: Path) -> None:
    path = tmp_path / "sales.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "销售数据"
    sheet.append(["产品", "销量"])
    for index in range(90):
        sheet.append([f"产品-{index}", index])
    workbook.save(path)

    result = parse_document(str(path))

    table_chunks = [chunk for chunk in result["chunks"] if chunk["metadata"]["kind"] == "table"]
    assert len(table_chunks) >= 3
    assert all(chunk["page"] == 1 for chunk in table_chunks)
    assert all(chunk["section_path"] == "销售数据" for chunk in table_chunks)
    assert all(chunk["metadata"]["sheet"] == "销售数据" for chunk in table_chunks)
    assert all("产品" in chunk["content"] for chunk in table_chunks)


def test_reports_stable_progress_stages(tmp_path: Path) -> None:
    path = tmp_path / "progress.docx"
    document = Document()
    document.add_paragraph("进度测试")
    document.save(path)
    events: list[tuple[str, int, str]] = []

    parse_document(
        str(path),
        progress=lambda stage, percent, message: events.append(
            (stage, percent, message)
        ),
    )

    assert [stage for stage, _, _ in events] == [
        "validating",
        "converting",
        "chunking",
        "finalizing",
    ]
    assert [percent for _, percent, _ in events] == [10, 45, 80, 95]


def test_rejects_xlsx_with_too_many_worksheet_members(tmp_path: Path) -> None:
    path = tmp_path / "oversized.xlsx"
    with zipfile.ZipFile(path, "w") as archive:
        for index in range(257):
            archive.writestr(
                f"xl/worksheets/sheet{index + 1}.xml",
                '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"/>',
            )

    with pytest.raises(DocumentParseError, match="工作表超过 256"):
        parse_document(str(path))


def test_rejects_xlsx_with_too_many_actual_cells(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    path = tmp_path / "too-many-cells.xlsx"
    monkeypatch.setattr(parser_module, "MAX_XLSX_CELLS", 2)
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "xl/worksheets/custom-data.xml",
            '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
            "<sheetData><row><c/><c/><c/></row></sheetData></worksheet>",
        )

    with pytest.raises(DocumentParseError, match="实际单元格数超过"):
        parse_document(str(path))


def test_rejects_pptx_with_too_many_slide_members(tmp_path: Path) -> None:
    path = tmp_path / "oversized.pptx"
    with zipfile.ZipFile(path, "w") as archive:
        for index in range(1_001):
            archive.writestr(
                f"ppt/slides/custom-{index + 1}.xml",
                '<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>',
            )

    with pytest.raises(DocumentParseError, match="幻灯片超过 1,000"):
        parse_document(str(path))


def test_cli_emits_jsonl_progress_and_result(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch, capsys: pytest.CaptureFixture[str]
) -> None:
    path = tmp_path / "protocol.docx"
    document = Document()
    document.add_paragraph("协议测试")
    document.save(path)
    monkeypatch.setattr(sys, "argv", ["document_parserd.py", "--path", str(path)])

    assert parser_main() == 0

    messages = [json.loads(line) for line in capsys.readouterr().out.splitlines()]
    assert [message["type"] for message in messages] == [
        "progress",
        "progress",
        "progress",
        "progress",
        "result",
    ]
    assert messages[-1]["payload"]["chunks"][0]["content"] == "协议测试"
